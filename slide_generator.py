import os 
from pptx import Presentation
from pptx.util import Inches 
from dotenv import load_dotenv


class SlideGenerator:
    def __init__(self, output_dir = "output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def create_slide_deck(self,title,bullet_points):
        prs = Presentation()

        # Title slides
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Content slides
        for i in range(0,len(bullet_points),5):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]

            title_shape.text = f"{title} (part {i//5 + 1})"

            for bullet in bullet_points[i:i+5]:
                p = content_shape.text_frame.add_paragraph()
                p.text = bullet
            
            # save presentation

            pptx_path = os.path.join(self.output_dir, "generated_presentation.pptx")
            prs.save(pptx_path)
            print(f"\n powerpoint slides saved to : {pptx_path}")

if __name__ == "__main__":
    sample_title = "AI-Based Slide Generation"
    sample_points = [
        "The AI Revolution refers to the rapid growth and integration of artificial intelligence technologies across industries and daily life.",

         "It has been fueled by advances in machine learning, deep learning, and neural networks, combined with the availability of massive amounts of data and powerful cloud computing resources.",

         "Natural language processing and computer vision breakthroughs have further accelerated the pace, enabling AI systems to understand, generate, and interact with humans in more natural ways.",

         "The impact of AI is visible across multiple domains.",

         "In healthcare, AI assists in diagnostics, drug discovery, and personalized treatments. In finance, it powers fraud detection, algorithmic trading, and customer support chatbots."
    ]

    slide_geneartor = SlideGenerator()
    slide_geneartor.create_slide_deck(sample_title,sample_points)